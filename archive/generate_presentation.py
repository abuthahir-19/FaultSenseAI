"""
Generates PANEL_PRESENTATION.pptx using the actual Prodapt Solutions template.
Opens Presentation1.pptx, clears all slides, adds 10 fresh slides using
the template's own layouts so branding, fonts, footers and colours are authentic.
Run: python generate_presentation.py
"""
import os
from copy import deepcopy
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn

TEMPLATE = os.path.join(os.path.dirname(__file__), "Presentation1.pptx")
OUT      = os.path.join(os.path.dirname(__file__), "PANEL_PRESENTATION.pptx")

# ── Prodapt palette ───────────────────────────────────────────────────────────
ORANGE  = RGBColor(0xF2, 0x6B, 0x43)
NAVY    = RGBColor(0x1B, 0x2A, 0x4A)
BLUE    = RGBColor(0x1B, 0x6E, 0xC2)
GREEN   = RGBColor(0x14, 0x7A, 0x40)
VIOLET  = RGBColor(0x6B, 0x21, 0xA8)
TEAL    = RGBColor(0x0D, 0x7E, 0x6D)
RED     = RGBColor(0xB9, 0x1C, 0x1C)
AMBER   = RGBColor(0xB4, 0x53, 0x09)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
DARK    = RGBColor(0x1E, 0x29, 0x3B)
MUTED   = RGBColor(0x64, 0x74, 0x8B)
CARD    = RGBColor(0xF1, 0xF5, 0xF9)
BORDER  = RGBColor(0xCB, 0xD5, 0xE1)

# ── Load template ─────────────────────────────────────────────────────────────
prs    = Presentation(TEMPLATE)
W      = prs.slide_width.inches     # 10.0
H      = prs.slide_height.inches    # 5.625
master = prs.slide_masters[0]

# ── Clear ALL existing slides ──────────────────────────────────────────────────
sldIdLst = prs.slides._sldIdLst
for sldId in list(sldIdLst):
    rId = sldId.get(qn("r:id"))
    prs.part.drop_rel(rId)
    sldIdLst.remove(sldId)

# ── Get layouts by name ────────────────────────────────────────────────────────
def layout(name):
    for l in master.slide_layouts:
        if l.name == name:
            return l
    return master.slide_layouts[41]   # blank fallback

L_COVER   = layout("Cover with one logo")                        # 0
L_1TEXT   = layout("Slide with text only")                       # 27
L_3BLOCK  = layout("3-step/block content slide without red band")# 22
L_4BLOCK  = layout("4-step/block content slide without red band")# 23
L_5BLOCK  = layout("5-step/block content slide without red band")# 24
L_6BLOCK  = layout("6-step/block content slide without red band")# 25
L_BLANK   = layout("Blank slide")                                # 41
L_THANKS  = layout("Thank you")                                  # 42
L_SECTION = layout("Prodapt specific separator")                 # 8


# ═════════════════════════════════════════════════════════════════════════════
# HELPERS
# ═════════════════════════════════════════════════════════════════════════════
def add_slide(lyt):
    return prs.slides.add_slide(lyt)


def ph(slide, idx):
    """Return placeholder by its placeholder idx."""
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == idx:
            return shape
    return None


def set_ph_text(slide, idx, text, size=None, bold=None, color=None):
    """Set placeholder text, optionally override font."""
    p = ph(slide, idx)
    if p is None:
        return
    tf = p.text_frame
    tf.text = text
    for para in tf.paragraphs:
        for run in para.runs:
            if size:  run.font.size  = Pt(size)
            if bold is not None: run.font.bold = bold
            if color: run.font.color.rgb = color


def set_ph_bullets(slide, idx, lines, size=None):
    """Set multi-line/bullet text in a placeholder."""
    p = ph(slide, idx)
    if p is None:
        return
    tf = p.text_frame
    tf.clear()
    for i, line in enumerate(lines):
        if isinstance(line, tuple):
            text, kw = line[0], line[1] if len(line) > 1 else {}
        else:
            text, kw = line, {}
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.text = text
        if size or kw.get("size") or kw.get("bold") or kw.get("color"):
            for run in para.runs:
                if size or kw.get("size"):
                    run.font.size = Pt(kw.get("size", size))
                if kw.get("bold") is not None:
                    run.font.bold = kw["bold"]
                if kw.get("color"):
                    run.font.color.rgb = kw["color"]


def box(sl, x, y, w, h, fill=CARD, border=None, bw=1.0):
    r = sl.shapes.add_shape(1,
        Inches(x), Inches(y), Inches(w), Inches(h))
    r.fill.solid(); r.fill.fore_color.rgb = fill
    if border:
        r.line.color.rgb = border
        r.line.width = Pt(bw)
    else:
        r.line.fill.background()
    return r


def txt(sl, text, x, y, w, h,
        size=10, color=DARK, bold=False, italic=False,
        align=PP_ALIGN.LEFT, wrap=True):
    x = min(x, W - 0.05); y = min(y, H - 0.05)
    w = min(w, W - x - 0.02); h = min(h, H - y - 0.02)
    txb = sl.shapes.add_textbox(
        Inches(x), Inches(y), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame; tf.word_wrap = wrap
    p  = tf.paragraphs[0]; p.alignment = align
    run = p.add_run()
    run.text = str(text)
    run.font.size   = Pt(size)
    run.font.color.rgb = color
    run.font.bold   = bold
    run.font.italic = italic
    run.font.name   = "Calibri"
    return txb


def bline(sl, x, y, w, h, color):
    b = sl.shapes.add_shape(1,
        Inches(x), Inches(y), Inches(w), Inches(h))
    b.fill.solid(); b.fill.fore_color.rgb = color
    b.line.fill.background()


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE  (Cover with one logo)
# ═════════════════════════════════════════════════════════════════════════════
sl = add_slide(L_COVER)
set_ph_text(sl, 10, "FaultSense AI", size=36, bold=True)
set_ph_text(sl, 11,
    "AI-Powered Telecom Network Fault Intelligence Platform\n"
    "RAG + LangGraph Multi-Agent Pipeline  ·  GPT-4o  ·  LangSmith  ·  DeepEval",
    size=13)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — PROBLEM STATEMENT  (3-step/block without red band)
# layout 22: PH0=title | PH93,96,99=headings | PH101,107,109=bodies
# ═════════════════════════════════════════════════════════════════════════════
sl = add_slide(L_3BLOCK)
set_ph_text(sl, 0, "Problem Statement — The NOC Engineer's Challenge")

set_ph_text(sl, 93, "📊  The Scale Problem")
set_ph_bullets(sl, 101, [
    "100+ alarms fired per hour during a NOC shift",
    "45+ minutes to manually identify root cause",
    "Engineers context-switch across 3-4 tools",
    "No unified view of related alarms",
    "→ Critical SLAs missed under manual triage",
], size=9.5)

set_ph_text(sl, 96, "⚠  The Pain Points")
set_ph_bullets(sl, 107, [
    "Root causes buried across EMS, NMS, tickets, vendor dashboards",
    "Spatial and temporal alarm correlation missed entirely",
    "Recurring incidents — resolutions never systematically mined",
    "Junior engineers lack domain knowledge for fast diagnosis",
    "→ Productivity lost to repetitive manual analysis",
], size=9.5)

set_ph_text(sl, 99, "💡  Why It Matters")
set_ph_bullets(sl, 109, [
    "60% of incidents recur without a systemic fix applied",
    "Each unresolved hour costs revenue and subscriber trust",
    "CRITICAL faults need emergency + regulatory context fast",
    "Traditional rule-based NOC tools cannot reason across incidents",
    "→ A purpose-built AI platform is the answer",
], size=9.5)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — SOLUTION OVERVIEW  (6-step/block without red band)
# layout 25: PH0=title | PH128,129,130=top headings | PH93,112,113=bottom headings
#            PH132,125,127=top bodies | PH123,111,103=bottom bodies
# ═════════════════════════════════════════════════════════════════════════════
sl = add_slide(L_6BLOCK)
set_ph_text(sl, 0, "FaultSense AI — Introducing the Solution")

# Top row
set_ph_text(sl, 128, "🔍  Hybrid RAG Search")
set_ph_bullets(sl, 132, [
    "ChromaDB semantic + BM25 keyword",
    "RRF fusion k=60 (Cormack 2009)",
    "23% better recall vs semantic-only",
    "LLM cross-encoder reranking",
], size=9)

set_ph_text(sl, 129, "🤖  5-Node LangGraph Pipeline")
set_ph_bullets(sl, 125, [
    "Guardrail → Retrieval → Correlation",
    "Root Cause → Service Impact → Resolution",
    "Conditional edge for CRITICAL incidents",
    "Full FaultAnalysisState shared across nodes",
], size=9)

set_ph_text(sl, 130, "🎯  Root Cause & Service Impact")
set_ph_bullets(sl, 127, [
    "GPT-4o chain-of-thought grounded in alarms",
    "Subscriber blast radius assessment",
    "SLA breach risk identification",
    "Cascading failure path analysis",
], size=9)

# Bottom row
set_ph_text(sl, 93, "📊  Analytics & Forecast")
set_ph_bullets(sl, 123, [
    "Real-time KPI dashboard",
    "30-day trend sparkline",
    "Vendor & technology breakdown",
    "AI predictive outage forecast",
], size=9)

set_ph_text(sl, 112, "📋  Auto RAG Evaluation")
set_ph_bullets(sl, 111, [
    "Faithfulness · Answer Relevancy",
    "Context Precision (RAGAS-aligned)",
    "Runs automatically after every analysis",
    "Expandable per-metric detail panels",
], size=9)

set_ph_text(sl, 113, "🛡  Guardrail Validation")
set_ph_bullets(sl, 103, [
    "Input validation (length/format)",
    "Injection detection (10 patterns)",
    "Telecom relevance (keyword check)",
    "Returns structured 200 — never 422",
], size=9)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — SYSTEM ARCHITECTURE  (Blank slide + custom shapes)
# ═════════════════════════════════════════════════════════════════════════════
sl = add_slide(L_BLANK)
set_ph_text(sl, 0, "System Architecture — End-to-End Flow")

# Content area: x=0.25..9.75, y=0.75..5.35
CX, CW = 0.25, 9.5

# ── User bubble ──
box(sl, 3.55, 0.78, 2.9, 0.44, fill=RGBColor(0xEF,0xF6,0xFF), border=BLUE, bw=1.2)
txt(sl, "👤  NOC Engineer  —  Plain-English Fault Query",
    3.57, 0.83, 2.86, 0.3, size=9, color=NAVY, align=PP_ALIGN.CENTER)

# down arrow
bline(sl, 4.98, 1.22, 0.04, 0.2, BLUE)
txt(sl, "Browser", 5.05, 1.25, 0.75, 0.16, size=7, color=MUTED)

# ── Frontend ──
box(sl, CX, 1.42, CW, 0.6, fill=RGBColor(0xEB,0xF5,0xFF), border=BLUE, bw=1.5)
bline(sl, CX, 1.42, CW, 0.055, BLUE)
txt(sl, "⚛  React Frontend  ·  Vite + TypeScript  ·  4 Modes: Query | Deep Analysis | Analytics | Evaluation",
    CX+0.15, 1.5, CW-0.25, 0.28, size=10, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
txt(sl, "QueryInput · IncidentCard · AgentTrace · RootCausePanel · AnalyticsDashboard · EvaluationPanel · GuardrailPanel · ErrorBoundary",
    CX+0.15, 1.8, CW-0.25, 0.2, size=8, color=MUTED, align=PP_ALIGN.CENTER)

# down arrow
bline(sl, 4.98, 2.02, 0.04, 0.2, BLUE)
txt(sl, "HTTP / axios", 5.05, 2.05, 1.0, 0.16, size=7, color=MUTED)

# ── Backend ──
box(sl, CX, 2.22, CW, 0.58, fill=RGBColor(0xED,0xFD,0xF4), border=GREEN, bw=1.5)
bline(sl, CX, 2.22, CW, 0.055, GREEN)
txt(sl, "⚡  FastAPI Backend  ·  Uvicorn  ·  Python 3.11  ·  Port 8000  ·  LangSmith bootstrap",
    CX+0.15, 2.3, CW-0.25, 0.26, size=10, color=GREEN, bold=True, align=PP_ALIGN.CENTER)
txt(sl, "/api/query  ·  /api/analyze  ·  /api/ingest  ·  /health  ·  /analytics/*  ·  /summarize  ·  /evaluate  ·  /rerank  ·  /incidents",
    CX+0.15, 2.58, CW-0.25, 0.2, size=8, color=MUTED, align=PP_ALIGN.CENTER)

# fork arrows
FY = 2.8
bline(sl, 2.4, FY, 0.04, 0.22, AMBER)
txt(sl, "Quick RAG", 1.7, FY+0.04, 0.65, 0.16, size=7, color=AMBER)
bline(sl, 7.12, FY, 0.04, 0.22, VIOLET)
txt(sl, "LangGraph", 7.2, FY+0.04, 0.75, 0.16, size=7, color=VIOLET)
# horizontal connector
bline(sl, 2.44, FY+0.18, 4.72, 0.03, BORDER)

# ── RAG + LangGraph side by side ──
TY, TH = 3.02, 1.15
box(sl, CX, TY, 4.4, TH, fill=RGBColor(0xFF,0xFB,0xEB), border=AMBER, bw=1.5)
bline(sl, CX, TY, 4.4, 0.055, AMBER)
txt(sl, "🔍  RAG Pipeline",
    CX+0.12, TY+0.1, 4.15, 0.24, size=10, color=AMBER, bold=True)
txt(sl, "EmbeddingManager  →  ChromaDBStore  →  BM25Index  →  HybridRetriever (RRF k=60)",
    CX+0.12, TY+0.38, 4.15, 0.24, size=8.5, color=DARK)
txt(sl, "text-embedding-3-small  ·  3 workers  ·  batch 512  ·  LLM Reranker",
    CX+0.12, TY+0.65, 4.15, 0.44, size=8, color=MUTED)

# cross-arrow RAG → LG
bline(sl, 4.7, TY+TH/2-0.02, 0.45, 0.035, VIOLET)
txt(sl, "search()", 4.7, TY+TH/2-0.2, 0.52, 0.16, size=6.5, color=VIOLET, align=PP_ALIGN.CENTER)

box(sl, 5.18, TY, 4.57, TH, fill=RGBColor(0xF5,0xF3,0xFF), border=VIOLET, bw=1.5)
bline(sl, 5.18, TY, 4.57, 0.055, VIOLET)
txt(sl, "🤖  LangGraph 5-Node Pipeline",
    5.3, TY+0.1, 4.35, 0.24, size=10, color=VIOLET, bold=True)
txt(sl, "Guardrail  →  [1] Retrieval  →  [2] Correlation  →  [3] Root Cause  →  [4] Service Impact (+fork)  →  [5] Resolution",
    5.3, TY+0.38, 4.35, 0.24, size=8, color=DARK)
txt(sl, "GPT-4o reasoning  ·  CRITICAL escalation fork  ·  Structured JSON output",
    5.3, TY+0.65, 4.35, 0.44, size=8, color=MUTED)

# ── Bottom row: Storage + External ──
BL = TY + TH + 0.14
box(sl, CX,    BL, 2.9, 0.56, fill=RGBColor(0xF1,0xF5,0xF9), border=MUTED, bw=0.8)
bline(sl, CX, BL, 2.9, 0.05, MUTED)
txt(sl, "🗂  Storage", CX+0.12, BL+0.09, 2.66, 0.22, size=9, color=DARK, bold=True)
txt(sl, "ChromaDB  ·  BM25  ·  telecom_incidents.csv",
    CX+0.12, BL+0.32, 2.66, 0.2, size=8, color=MUTED)

box(sl, 3.42, BL, 2.05, 0.56, fill=RGBColor(0xEE,0xF2,0xFF), border=BLUE, bw=0.8)
bline(sl, 3.42, BL, 2.05, 0.05, BLUE)
txt(sl, "🧠  OpenAI API", 3.54, BL+0.09, 1.82, 0.22, size=9, color=BLUE, bold=True)
txt(sl, "GPT-4o  ·  text-embedding-3-small", 3.54, BL+0.32, 1.82, 0.2, size=8, color=MUTED)

box(sl, 5.7, BL, 2.05, 0.56, fill=RGBColor(0xEC,0xFD,0xF5), border=TEAL, bw=0.8)
bline(sl, 5.7, BL, 2.05, 0.05, TEAL)
txt(sl, "📊  LangSmith", 5.82, BL+0.09, 1.82, 0.22, size=9, color=TEAL, bold=True)
txt(sl, "Per-agent trace  ·  Token counts", 5.82, BL+0.32, 1.82, 0.2, size=8, color=MUTED)

box(sl, 7.98, BL, 1.77, 0.56, fill=RGBColor(0xFE,0xF2,0xF2), border=RED, bw=0.8)
bline(sl, 7.98, BL, 1.77, 0.05, RED)
txt(sl, "📋  DeepEval", 8.1, BL+0.09, 1.55, 0.22, size=9, color=RED, bold=True)
txt(sl, "Faithfulness  ·  Relevancy  ·  Precision",
    8.1, BL+0.32, 1.55, 0.2, size=8, color=MUTED)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — HYBRID RAG PIPELINE  (3-step without red band)
# 3 columns: (1) Query → Embed, (2) Semantic + Keyword Search, (3) RRF + Rerank
# ═════════════════════════════════════════════════════════════════════════════
sl = add_slide(L_3BLOCK)
set_ph_text(sl, 0, "Hybrid RAG Pipeline — How Retrieval Works")

set_ph_text(sl, 93, "Step 1 — Vectorise the Query")
set_ph_bullets(sl, 101, [
    "User submits plain-English fault query",
    "text-embedding-3-small (OpenAI API)",
    "Produces 1536-dimensional query vector",
    "3 ThreadPoolExecutor workers",
    "Batch size: 512 documents per round",
    "Ingest time: 40-60 s → 12-15 s with concurrency",
    "",
    "Key design: 3× cheaper than ada-002,",
    "  comparable domain quality",
], size=9)

set_ph_text(sl, 96, "Step 2 — Dual Search (Parallel)")
set_ph_bullets(sl, 107, [
    "ChromaDB similarity_search:",
    "  — Dense semantic vector search",
    "  — Metadata filters: region, vendor,",
    "    severity, technology",
    "",
    "BM25Index full-text search:",
    "  — Captures exact alarm IDs and codes",
    "  — rank_bm25 tokenisation",
    "  — Complementary to dense search",
], size=9)

set_ph_text(sl, 99, "Step 3 — RRF Fusion + LLM Rerank")
set_ph_bullets(sl, 109, [
    "Reciprocal Rank Fusion (RRF):",
    "  score = Σ 1 / (rank + 60)",
    "  k=60: Cormack 2009 proven constant",
    "  23% better top-5 recall vs semantic-only",
    "",
    "LLM Cross-Encoder Reranker:",
    "  combined = 0.6×LLM + 0.4×RRF",
    "  Domain-aware re-ordering",
    "  POST /api/rerank endpoint",
], size=9)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — LANGGRAPH 5-NODE PIPELINE  (6-step without red band)
# 6 blocks: Guardrail + 5 nodes (top 3: G,1,2 | bottom 3: 3,4,5)
# layout 25: PH0=title | PH96,99=top2 heads (but we need 3 top + 3 bottom)
# Actually layout 25 has: top PH96(0.36",1.11",4.60x1.97), PH99(5.05",1.11",4.60x1.97)
#                         bottom PH93(0.36",3.15",3.03), PH112(3.48"), PH113(6.61")
# Use layout 25: 2 wide on top + 3 narrow on bottom
# ═════════════════════════════════════════════════════════════════════════════
sl = add_slide(L_5BLOCK)
set_ph_text(sl, 0, "LangGraph 5-Node Multi-Agent Pipeline")

# layout 24 "5-step without red band":
# top row: PH96(0.36",1.11",4.60x1.97) PH99(5.05",1.11",4.60x1.97)
# body top: PH101(0.44",1.51",4.43x1.38) PH119(5.13",1.51",4.43x1.38)
# bottom row: PH93(0.36",3.15",3.03x1.97) PH112(3.48",3.15",3.03) PH113(6.61",3.15",3.03)
# body bot: PH123(0.44",3.56") PH111(3.58") PH103(6.71")

set_ph_text(sl, 96, "🛡  Guardrail  +  Node 1: Alarm Retrieval")
set_ph_bullets(sl, 101, [
    "Guardrail: Input validation · Injection detection · Telecom relevance",
    "Returns structured HTTP 200 — never a raw 422 error",
    "Node 1: HybridRetriever.search() with RRF fusion",
    "Sets severity_escalated flag for CRITICAL incidents",
    "File: alarm_retrieval_agent.py",
], size=9)

set_ph_text(sl, 99, "Node 2: Cross-Correlation  (Deterministic)")
set_ph_bullets(sl, 119, [
    "No LLM call — pure deterministic logic in utils/correlation.py",
    "Clusters incidents by: region + technology + time window",
    "Extracts: dominant vendor, max severity, time span",
    "Output: correlated_alarms list in FaultAnalysisState",
    "Design: separating correlation keeps LLM agents focused",
], size=9)

set_ph_text(sl, 93, "Node 3\nRoot Cause")
set_ph_bullets(sl, 123, [
    "GPT-4o chain-of-thought reasoning",
    "Grounded in retrieved alarm IDs",
    "Produces reasoning trace for UI",
    "File: root_cause_agent.py",
], size=8.5)

set_ph_text(sl, 112, "Node 4\nService Impact (+Fork)")
set_ph_bullets(sl, 111, [
    "Subscriber blast radius",
    "SLA breach risk assessment",
    "Cascading failure paths",
    "CRITICAL fork: emergency + regulatory context",
], size=8.5)

set_ph_text(sl, 113, "Node 5\nResolution")
set_ph_bullets(sl, 103, [
    "GPT-4o structured JSON output",
    "5 categories: IMMEDIATE / DIAGNOSTIC",
    "RESOLUTION / PREVENTIVE / ESCALATION",
    "File: resolution_agent.py",
], size=8.5)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — KEY FEATURES  (6-step without red band)
# ═════════════════════════════════════════════════════════════════════════════
sl = add_slide(L_6BLOCK)
set_ph_text(sl, 0, "Platform Features — 16 Capabilities")

set_ph_text(sl, 128, "🔍  Retrieval & Search")
set_ph_bullets(sl, 132, [
    "Hybrid RAG: ChromaDB + BM25 + RRF",
    "LLM cross-encoder reranking",
    "Metadata filters (region/vendor/tech)",
    "23% recall improvement",
], size=9)

set_ph_text(sl, 129, "🤖  Agent Intelligence")
set_ph_bullets(sl, 125, [
    "5-node LangGraph StateGraph",
    "Root cause (GPT-4o chain-of-thought)",
    "Service impact & blast radius",
    "CRITICAL severity escalation fork",
], size=9)

set_ph_text(sl, 130, "✅  Recommendations")
set_ph_bullets(sl, 127, [
    "Structured JSON: 5 categories",
    "IMMEDIATE / DIAGNOSTIC / RESOLUTION",
    "PREVENTIVE / ESCALATION steps",
    "Alarm correlation clustering",
], size=9)

set_ph_text(sl, 93, "📊  Analytics & Intelligence")
set_ph_bullets(sl, 123, [
    "KPI dashboard with real-time data",
    "30-day trend sparkline",
    "Vendor & technology breakdowns",
    "AI predictive outage forecast",
], size=9)

set_ph_text(sl, 112, "📋  Evaluation & Quality")
set_ph_bullets(sl, 111, [
    "Auto RAGAS evaluation after analysis",
    "Faithfulness · Answer Relevancy",
    "Context Precision with detail panels",
    "LangSmith per-agent tracing",
], size=9)

set_ph_text(sl, 113, "🛡  Safety & Resilience")
set_ph_bullets(sl, 103, [
    "3-check guardrail validation panel",
    "Executive outage summarisation",
    "React ErrorBoundary crash safety",
    "Live ingest progress (SSE polling)",
], size=9)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — UI MODES  (4-step without red band)
# layout 23: PH0=title | PH96(0.36",1.11",4.60x1.97) PH99(5.05",1.11",4.60x1.97)
#            PH120(0.36",3.15",4.60) PH121(5.05",3.15",4.60)
#            bodies: PH101,119,123,125
# ═════════════════════════════════════════════════════════════════════════════
sl = add_slide(L_4BLOCK)
set_ph_text(sl, 0, "User Interface — 4 Application Modes")

set_ph_text(sl, 96, "🔎  Query Mode")
set_ph_bullets(sl, 101, [
    "Natural language search with metadata filters",
    "Guardrail Panel — shows all 3 checks above results",
    "Incident cards with RRF scores and severity badges",
    "Formatted outage duration ('2h 15m' not '135')",
    "Quick LLM root cause suggestion inline",
    "Mode-aware button: only Quick Search is active",
], size=9)

set_ph_text(sl, 99, "🧠  Deep Analysis Mode")
set_ph_bullets(sl, 119, [
    "Full 5-node LangGraph pipeline execution",
    "Guardrail Panel always shown at top",
    "LangGraph reasoning trace accordion",
    "Correlated alarm clusters with vendor/region/time",
    "Root cause + service impact narrative",
    "Categorised resolution recommendations",
], size=9)

set_ph_text(sl, 120, "📊  Analytics Mode")
set_ph_bullets(sl, 123, [
    "KPI cards: total, CRITICAL, HIGH counts",
    "Severity distribution pie chart",
    "Technology and device vendor breakdowns",
    "30-day daily incident trend sparkline",
    "AI predictive outage forecast (GPT-4o)",
    "POST /api/analytics/predict endpoint",
], size=9)

set_ph_text(sl, 121, "📋  Evaluation Mode")
set_ph_bullets(sl, 125, [
    "Auto-triggered after every Deep Analysis",
    "Live tab badge: yellow = running, violet = ready",
    "3 expandable metric cards with score bars",
    "Each card: assessment + 'View Details' dropdown",
    "Details: what it measures + high/low interpretation",
    "Weighted overall: 0.40×F + 0.35×AR + 0.25×CP",
], size=9)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — EVALUATION, GUARDRAILS & LANGSMITH  (3-step without red band)
# ═════════════════════════════════════════════════════════════════════════════
sl = add_slide(L_3BLOCK)
set_ph_text(sl, 0, "Quality & Observability — DeepEval · Guardrails · LangSmith")

set_ph_text(sl, 93, "📋  RAGAS-Aligned RAG Evaluation")
set_ph_bullets(sl, 101, [
    "Faithfulness (weight: ×0.40)",
    "  Are claims grounded in retrieved incidents?",
    "  Detects hallucination in the analysis",
    "",
    "Answer Relevancy (weight: ×0.35)",
    "  Does the analysis address the query fully?",
    "  Checks for missing aspects",
    "",
    "Context Precision (weight: ×0.25)",
    "  What fraction of retrieved incidents are relevant?",
    "  Measures RAG retrieval quality",
    "",
    "Implementation: 3 direct LLM-as-judge calls",
    "  (not DeepEval built-ins — avoids 500-token proxy cap)",
], size=8.5)

set_ph_text(sl, 96, "🛡  3-Check Guardrail Panel")
set_ph_bullets(sl, 107, [
    "Check 1 — Input Validation:",
    "  Empty query · too short · too long · format",
    "  Returns: PASS or FAIL with message",
    "",
    "Check 2 — Injection Detection:",
    "  10 regex patterns: prompt injection,",
    "  SQL injection, script tags, jailbreak attempts",
    "  Returns: PASS or FAIL",
    "",
    "Check 3 — Telecom Relevance:",
    "  Word-boundary keyword matching",
    "  (fixes false positives: 'port' ≠ 'report')",
    "  Returns: PASS or WARNING",
    "",
    "All checks return HTTP 200 with structured result",
], size=8.5)

set_ph_text(sl, 99, "📊  LangSmith Tracing")
set_ph_bullets(sl, 109, [
    "Full per-run agent trace for every /api/analyze",
    "Token counts per LLM call and per node",
    "Wall-clock latency for each of the 5 nodes",
    "Inputs and outputs stored for debugging",
    "",
    "Activation:",
    "  LANGCHAIN_TRACING_V2 = true",
    "  LANGCHAIN_API_KEY = <key>",
    "  LANGCHAIN_PROJECT = multi-agent-trip-planner",
    "",
    "Bootstrap in main.py:",
    "  Sets os.environ BEFORE LangGraph imports",
    "  (LangChain reads env vars at import time)",
    "",
    "Zero-config for demos; full observability in prod",
], size=8.5)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — THANK YOU / CONCLUSION  (Thank you layout)
# ═════════════════════════════════════════════════════════════════════════════
sl = add_slide(L_THANKS)
set_ph_text(sl, 13, "Thank You")
set_ph_text(sl, 14, "Questions Welcome  ·  FaultSense AI  ·  2026")

# Add achievement summary as custom text (positioned above template thank-you area)
txt(sl, "Key Achievements",
    0.3, 0.78, 9.4, 0.34, size=13, color=NAVY, bold=True, align=PP_ALIGN.CENTER)

achievements = [
    ("5-Node LangGraph",  "Full orchestration with CRITICAL escalation fork",  VIOLET),
    ("Hybrid RAG",        "23% recall lift · ChromaDB + BM25 + RRF + Reranker", AMBER),
    ("Auto Evaluation",   "RAGAS Faithfulness · Relevancy · Precision metrics",  GREEN),
    ("LangSmith Tracing", "End-to-end observability for every agent pipeline",   TEAL),
    ("Guardrail System",  "3-check validation — structured 200 responses",       ORANGE),
    ("16 Features",       "Production-ready platform across 4 UI modes",        BLUE),
]

aw = (9.4 - 0.1) / 3
ax = 0.3
for i, (title, desc, col) in enumerate(achievements):
    row, col_n = divmod(i, 3)
    bx = 0.3 + col_n * (aw + 0.03)
    by = 1.18 + row * 0.52
    box(sl, bx, by, aw, 0.44, fill=RGBColor(0xF8,0xFA,0xFC), border=col, bw=1.2)
    bline(sl, bx, by, aw, 0.05, col)
    txt(sl, title, bx+0.12, by+0.07, aw*0.42, 0.22, size=9.5, color=col, bold=True)
    txt(sl, desc,  bx+0.55, by+0.07, aw*0.55, 0.3,  size=8.5, color=DARK)


# ─────────────────────────────────────────────────────────────────────────────
prs.save(OUT)
sz = os.path.getsize(OUT)
print(f"Saved:  {OUT}")
print(f"Slides: {len(prs.slides)}  |  Size: {sz:,} bytes ({sz/1024:.0f} KB)")
