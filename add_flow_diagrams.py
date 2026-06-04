"""
Adds 3 flow-diagram slides to PANEL_PRESENTATION.pptx using the exact
colour palette from the System Architecture slide (Slide 4).

Slides added:
  11 — Data Ingestion Pipeline
  12 — Quick Search Flow
  13 — Deep Analysis Flow

Run: python add_flow_diagrams.py
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

PPTX = os.path.join(os.path.dirname(__file__), "PANEL_PRESENTATION.pptx")

# ── Exact palette from the architecture slide ─────────────────────────────────
ORANGE  = RGBColor(0xF2, 0x6B, 0x43)   # Prodapt brand / Guardrail
NAVY    = RGBColor(0x1B, 0x2A, 0x4A)   # slide header / User node
BLUE    = RGBColor(0x1B, 0x6E, 0xC2)   # Frontend / Search step
GREEN   = RGBColor(0x14, 0x7A, 0x40)   # Backend / Output/Results step
VIOLET  = RGBColor(0x6B, 0x21, 0xA8)   # LangGraph / AI Analysis step
TEAL    = RGBColor(0x0D, 0x7E, 0x6D)   # LangSmith / BM25 step
AMBER   = RGBColor(0xB4, 0x53, 0x09)   # RAG / Embeddings step
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
DARK    = RGBColor(0x1E, 0x29, 0x3B)
MUTED   = RGBColor(0x64, 0x74, 0x8B)
BORDER  = RGBColor(0xCB, 0xD5, 0xE1)

# Light fills matching the architecture card fills
_LT = {
    NAVY:   RGBColor(0xF1, 0xF5, 0xF9),
    ORANGE: RGBColor(0xFF, 0xF7, 0xED),
    BLUE:   RGBColor(0xEB, 0xF5, 0xFF),
    GREEN:  RGBColor(0xED, 0xFD, 0xF4),
    VIOLET: RGBColor(0xF5, 0xF3, 0xFF),
    TEAL:   RGBColor(0xEC, 0xFD, 0xF5),
    AMBER:  RGBColor(0xFF, 0xFB, 0xEB),
}

prs    = Presentation(PPTX)
master = prs.slide_masters[0]
W      = prs.slide_width.inches    # 10.0
H      = prs.slide_height.inches   # 5.625

def get_layout(name):
    for l in master.slide_layouts:
        if l.name == name:
            return l
    return master.slide_layouts[41]

L_BLANK = get_layout("Blank slide")


# ═════════════════════════════════════════════════════════════════════════════
# DRAWING PRIMITIVES
# ═════════════════════════════════════════════════════════════════════════════
def solid_rect(sl, x, y, w, h, fill, border=None, bw=1.0, rounded=False):
    shp = sl.shapes.add_shape(
        5 if rounded else 1,
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if border:
        shp.line.color.rgb = border
        shp.line.width = Pt(bw)
    else:
        shp.line.fill.background()
    return shp


def oval(sl, cx, cy, r, fill, border=None, bw=1.0):
    shp = sl.shapes.add_shape(
        9,   # MSO_AUTO_SHAPE_TYPE.OVAL
        Inches(cx - r), Inches(cy - r),
        Inches(r * 2),  Inches(r * 2)
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if border:
        shp.line.color.rgb = border
        shp.line.width = Pt(bw)
    else:
        shp.line.fill.background()
    return shp


def txb(sl, text, x, y, w, h,
        size=10, color=DARK, bold=False, italic=False,
        align=PP_ALIGN.CENTER, wrap=True):
    x = min(x, W - 0.05); y = min(y, H - 0.05)
    w = min(w, W - x - 0.02); h = min(h, H - y - 0.02)
    tb = sl.shapes.add_textbox(
        Inches(x), Inches(y), Inches(w), Inches(h))
    tb.word_wrap = wrap
    tf = tb.text_frame; tf.word_wrap = wrap
    p  = tf.paragraphs[0]; p.alignment = align
    run = p.add_run()
    run.text = str(text)
    run.font.size   = Pt(size)
    run.font.color.rgb = color
    run.font.bold   = bold
    run.font.italic = italic
    run.font.name   = "Calibri"
    return tb


def set_ph_title(sl, text, size=16):
    for shp in sl.placeholders:
        if shp.placeholder_format.idx == 0:
            shp.text = text
            for para in shp.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(size)
                    run.font.bold = True
            return


# ═════════════════════════════════════════════════════════════════════════════
# FLOW DIAGRAM BUILDER
# ═════════════════════════════════════════════════════════════════════════════
def draw_flow(sl, steps):
    """
    steps: list of dicts  {emoji, title, sub, lines[2], accent}
    Draws the complete horizontal flow diagram on the slide.
    Uses exact Prodapt palette from the architecture diagram.
    """
    N = len(steps)

    # ── Layout geometry ───────────────────────────────────────────────────
    CY   = 0.90   # card top
    CH   = 2.75   # card height
    if N == 5:
        CW, AW = 1.52, 0.42
    else:                       # N == 6
        CW, AW = 1.27, 0.37
    TOTAL = N * CW + (N - 1) * AW
    CX0   = (W - TOTAL) / 2     # first card left edge
    ICY   = CY + 0.70           # icon circle centre y
    IR    = 0.30                 # icon circle radius

    # ── Dark header bar (matches Slide 4 tier headers) ────────────────────
    hdr_y = CY - 0.02
    solid_rect(sl, CX0 - 0.08, hdr_y - 0.01,
               TOTAL + 0.16, 0.055, NAVY)

    for i, step in enumerate(steps):
        cx = CX0 + i * (CW + AW)   # card left edge
        acc   = step["accent"]
        light = _LT.get(acc, RGBColor(0xF8,0xFA,0xFC))

        # ── Card background ───────────────────────────────────────────────
        solid_rect(sl, cx, CY, CW, CH,
                   fill=light, border=acc, bw=1.4, rounded=True)

        # ── Accent top strip ──────────────────────────────────────────────
        solid_rect(sl, cx, CY, CW, 0.065, fill=acc, rounded=False)

        # ── Step number badge ─────────────────────────────────────────────
        NR = 0.145
        oval(sl, cx + 0.21, CY + 0.25, NR, fill=acc)
        txb(sl, str(i + 1),
            cx + 0.21 - NR, CY + 0.25 - NR, NR * 2, NR * 2,
            size=9.5, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

        # ── Icon outer glow ring ──────────────────────────────────────────
        icx = cx + CW / 2
        oval(sl, icx, ICY, IR + 0.08, fill=acc,
             border=None)
        # Make it semi-transparent via a slightly lighter shade
        # (python-pptx doesn't support opacity on fills easily, so use ring colour)
        ring_shp = sl.shapes[-1]
        ring_shp.fill.fore_color.rgb = _LT.get(acc, light)

        # ── Icon circle ───────────────────────────────────────────────────
        oval(sl, icx, ICY, IR, fill=light, border=acc, bw=1.8)

        # ── Emoji ─────────────────────────────────────────────────────────
        txb(sl, step["emoji"],
            icx - IR, ICY - IR, IR * 2, IR * 2,
            size=18, color=DARK, align=PP_ALIGN.CENTER)

        # ── Title ─────────────────────────────────────────────────────────
        ty = CY + 1.18
        txb(sl, step["title"],
            cx + 0.07, ty, CW - 0.14, 0.28,
            size=10.5, color=DARK, bold=True, align=PP_ALIGN.CENTER)

        # ── Sub-title ─────────────────────────────────────────────────────
        txb(sl, step["sub"],
            cx + 0.07, ty + 0.29, CW - 0.14, 0.22,
            size=8.5, color=acc, bold=True, align=PP_ALIGN.CENTER)

        # ── Divider ───────────────────────────────────────────────────────
        solid_rect(sl, cx + 0.12, ty + 0.54, CW - 0.24, 0.012,
                   fill=BORDER, rounded=False)

        # ── Detail lines ──────────────────────────────────────────────────
        for j, line in enumerate(step.get("lines", [])):
            ly = ty + 0.62 + j * 0.22
            if ly + 0.2 > CY + CH:
                break
            txb(sl, line,
                cx + 0.07, ly, CW - 0.14, 0.2,
                size=8, color=MUTED, align=PP_ALIGN.CENTER)

        # ── Arrow to next step ────────────────────────────────────────────
        if i < N - 1:
            next_acc = steps[i + 1]["accent"]
            ax       = cx + CW + 0.04
            ay       = CY + CH / 2 - 0.015
            # Arrow shaft
            solid_rect(sl, ax, ay, AW - 0.1, 0.028,
                       fill=next_acc, rounded=False)
            # Arrowhead triangle (approximated as narrow diamond)
            tip_x = ax + AW - 0.1
            tri = sl.shapes.add_shape(
                13,  # RIGHT_TRIANGLE as rough arrowhead
                Inches(tip_x - 0.005), Inches(ay - 0.065),
                Inches(0.105), Inches(0.16)
            )
            tri.fill.solid(); tri.fill.fore_color.rgb = next_acc
            tri.line.fill.background()

            # Arrow label
            lbl = step.get("arrow_label", "")
            if lbl:
                lbl_x = ax + (AW - 0.1) / 2 - 0.18
                lbl_y = ay - 0.23
                solid_rect(sl, lbl_x, lbl_y, 0.36, 0.19,
                           fill=WHITE, border=next_acc, bw=0.8, rounded=True)
                txb(sl, lbl, lbl_x, lbl_y, 0.36, 0.19,
                    size=7.5, color=next_acc, bold=True, align=PP_ALIGN.CENTER)

    # ── Bottom tech bar ───────────────────────────────────────────────────
    bar_y = CY + CH + 0.18
    solid_rect(sl, CX0 - 0.08, bar_y, TOTAL + 0.16, 0.33,
               fill=NAVY, rounded=False)
    cw_each = TOTAL / N
    for i, step in enumerate(steps):
        lx = CX0 + i * cw_each + cw_each / 2
        tech = step.get("tech", "")
        txb(sl, tech,
            CX0 + i * cw_each, bar_y, cw_each, 0.33,
            size=7.5, color=RGBColor(0x94, 0xA3, 0xB8),
            italic=True, align=PP_ALIGN.CENTER)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE DEFINITIONS
# ═════════════════════════════════════════════════════════════════════════════

# ── SLIDE 11 — INGESTION DATA FLOW ───────────────────────────────────────────
ingestion_steps = [
    {
        "emoji": "📄",
        "title": "CSV File",
        "sub":   "Source Data",
        "lines": ["telecom_incidents.csv", "9,827 records"],
        "accent": GREEN,
        "tech":   "pandas · data loader",
        "arrow_label": "load",
    },
    {
        "emoji": "⚡",
        "title": "Ingestion Pipeline",
        "sub":   "Backend Trigger",
        "lines": ["POST /api/ingest", "Background worker"],
        "accent": BLUE,
        "tech":   "FastAPI BackgroundTask",
        "arrow_label": "embed",
    },
    {
        "emoji": "🧠",
        "title": "Generate Embeddings",
        "sub":   "text-embedding-3-small",
        "lines": ["OpenAI API · 3 workers", "Batch 512 · 12–15 sec"],
        "accent": AMBER,
        "tech":   "OpenAI text-embedding-3-small",
        "arrow_label": "store",
    },
    {
        "emoji": "🗄",
        "title": "Store in ChromaDB",
        "sub":   "Vector Database",
        "lines": ["Persistent local store", "1536-dim vectors"],
        "accent": VIOLET,
        "tech":   "ChromaDB persistent",
        "arrow_label": "index",
    },
    {
        "emoji": "🔍",
        "title": "Build BM25 Index",
        "sub":   "Keyword Search Index",
        "lines": ["rank_bm25 library", "In-memory · ready to query"],
        "accent": TEAL,
        "tech":   "rank_bm25 · in-memory",
        "arrow_label": "",
    },
]

# ── SLIDE 12 — QUICK SEARCH FLOW ─────────────────────────────────────────────
quicksearch_steps = [
    {
        "emoji": "🔎",
        "title": "User Query",
        "sub":   "Plain English Input",
        "lines": ["Type fault description", "No technical knowledge"],
        "accent": NAVY,
        "tech":   "React QueryInput",
        "arrow_label": "check",
    },
    {
        "emoji": "🛡",
        "title": "Guardrail Check",
        "sub":   "Safety & Relevance",
        "lines": ["3 validation checks", "Safe · telecom-related?"],
        "accent": ORANGE,
        "tech":   "utils/guardrails.py",
        "arrow_label": "search",
    },
    {
        "emoji": "📚",
        "title": "Hybrid Search",
        "sub":   "Find Similar Incidents",
        "lines": ["ChromaDB + BM25 + RRF", "Best matches ranked"],
        "accent": BLUE,
        "tech":   "HybridRetriever · RRF k=60",
        "arrow_label": "suggest",
    },
    {
        "emoji": "💡",
        "title": "Quick Root Cause",
        "sub":   "AI Suggestion",
        "lines": ["Reads similar incidents", "2–3 sentence insight"],
        "accent": AMBER,
        "tech":   "GPT-4o · /api/query",
        "arrow_label": "show",
    },
    {
        "emoji": "📋",
        "title": "Results Shown",
        "sub":   "Incident Cards",
        "lines": ["Ranked by RRF score", "Severity · duration"],
        "accent": GREEN,
        "tech":   "React IncidentCard",
        "arrow_label": "",
    },
]

# ── SLIDE 13 — DEEP ANALYSIS FLOW ────────────────────────────────────────────
deepanalysis_steps = [
    {
        "emoji": "🔎",
        "title": "User Query",
        "sub":   "Plain English Input",
        "lines": ["Deep investigation", "Full pipeline triggered"],
        "accent": NAVY,
        "tech":   "React QueryInput",
        "arrow_label": "check",
    },
    {
        "emoji": "🛡",
        "title": "Guardrail Check",
        "sub":   "Safety & Relevance",
        "lines": ["3 checks run", "Returns structured result"],
        "accent": ORANGE,
        "tech":   "utils/guardrails.py",
        "arrow_label": "retrieve",
    },
    {
        "emoji": "📚",
        "title": "Retrieve Incidents",
        "sub":   "Find Related Alarms",
        "lines": ["Hybrid RAG search", "Top-K ranked incidents"],
        "accent": BLUE,
        "tech":   "HybridRetriever · RRF",
        "arrow_label": "correlate",
    },
    {
        "emoji": "🔗",
        "title": "Correlate Alarms",
        "sub":   "Find Patterns",
        "lines": ["Group by region + tech", "Identify affected vendors"],
        "accent": AMBER,
        "tech":   "utils/correlation.py",
        "arrow_label": "analyse",
    },
    {
        "emoji": "🎯",
        "title": "AI Analysis",
        "sub":   "Root Cause + Impact",
        "lines": ["GPT-4o explains why", "Assesses service impact"],
        "accent": VIOLET,
        "tech":   "LangGraph nodes 3 & 4",
        "arrow_label": "generate",
    },
    {
        "emoji": "✅",
        "title": "Recommendations",
        "sub":   "Fix Guide + Evaluation",
        "lines": ["Step-by-step actions", "Auto quality score"],
        "accent": GREEN,
        "tech":   "Node 5 · DeepEval",
        "arrow_label": "",
    },
]

# ═════════════════════════════════════════════════════════════════════════════
# ADD SLIDES
# ═════════════════════════════════════════════════════════════════════════════
slide_defs = [
    ("Data Ingestion Pipeline",
     "How raw CSV data becomes a searchable vector knowledge base",
     ingestion_steps),
    ("Quick Search — High-Level Flow",
     "Hybrid retrieval with an inline AI root cause hint — results in seconds",
     quicksearch_steps),
    ("Deep Analysis — AI Agent Pipeline",
     "Full 5-node LangGraph pipeline from alarm to root cause, service impact, and recommendations",
     deepanalysis_steps),
]

start_idx = len(prs.slides) + 1
for title, subtitle, steps in slide_defs:
    sl = prs.slides.add_slide(L_BLANK)
    # Title placeholder
    set_ph_title(sl, title, size=16)
    # Subtitle (custom text below title)
    txb(sl, subtitle,
        0.25, 0.72, 9.5, 0.22,
        size=9, color=MUTED, italic=True, align=PP_ALIGN.CENTER)
    # Thin orange Prodapt underline
    solid_rect(sl, 0.0, 0.65, W, 0.04, fill=ORANGE)
    # Draw the flow diagram
    draw_flow(sl, steps)

prs.save(PPTX)
total = len(prs.slides)
print(f"Saved: {PPTX}")
print(f"Total slides now: {total}  (added slides {start_idx}–{total})")
